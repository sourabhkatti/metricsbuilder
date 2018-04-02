import argparse
import base64
import calendar
import configparser
import csv
import json
import os
import re
import sys
from datetime import datetime, timedelta

import requests


class controller:
    ####################################################################
    #  Location of properties with connection details to the Contrast UI
    properties_file = "script.properties"

    header = {}
    application_licensed_status = {}

    endTimeEpoch = 0
    startTimeEpoch = 0
    header = {}

    groups = {}
    group_details = {}

    def __init__(self):
        parser = argparse.ArgumentParser(description='Communicate with the Contrast Rest API')

        parser.add_argument("-c", help="Specify the path to the scripts.properties file", nargs=1, type=str, metavar="")
        parser.add_argument("-o", help='Specify the path to write text files to', nargs=1,
                            type=str, metavar="")
        parser.add_argument("-i", help='Specify the organization ID', nargs=1,
                            type=str, metavar="")
        parser.add_argument("-t", help='Specify the teamserver URL', nargs=1,
                            type=str, metavar="")
        parser.add_argument("-a", help='Specify the API Key', nargs=1,
                            type=str, metavar="")
        parser.add_argument("-s", help='Specify the Service Key', nargs=1,
                            type=str, metavar="")
        parser.add_argument("-u", help='Specify the Username', nargs=1,
                            type=str, metavar="")
        parser.add_argument("--LibraryMetrics", help="Retrieve library metrics", action='store_true')
        parser.add_argument("--VulnerabilityMetrics", help='Retrieve trending vulnerability metrics',
                            action='store_true')
        parser.add_argument("--ApplicationMetrics", help="Retrieve application metrics", action='store_true')

        argt = parser.parse_args()

        if argt.o:
            self.outputpath = argt.o[0]
        else:
            self.outputpath = os.getcwd()
        if argt.c:
            self.properties_file = argt.c[0]
        else:
            self.properties_file = "script.properties"

        self.getScriptConfiguration()
        print("Writing output to", self.outputpath)

        if argt.VulnerabilityMetrics:
            self.VulnerabilityTrendMetrics(application_metrics=True)
        if argt.ApplicationMetrics:
            self.ApplicationMetrics()
        if argt.LibraryMetrics:
            self.LibraryMetrics()

        if "LibraryMetrics" in sys.argv:
            self.LibraryMetrics()
        if "ApplicationMetrics" in sys.argv:
            self.ApplicationMetrics()
        if "VulnerabilityMetrics" in sys.argv:
            self.VulnerabilityTrendMetrics(application_metrics=True)

    # noinspection PyShadowingBuiltins
    def getScriptConfiguration(self):
        """
        Open properties file and set the Contrast UI connection details
        """
        cfp = configparser.ConfigParser()
        cfp.read(self.properties_file)

        # Get all properties
        try:
            self.USERNAME = cfp.get("Contrast UI Details", "contrast.username")
        except:
            print("Please specify contrast.username in the script.properties")
        try:
            self.ORGANIZATION_ID = cfp.get("Contrast UI Details", "contrast.organization.id")
        except:
            print("Please specify contrast.organizationi.id in the script.properties")
        try:
            self.API_KEY = cfp.get("Contrast UI Details", "contrast.apikey")
        except:
            print("Please specify contrast.apikey in the script.properties")
        try:
            self.SERVICE_KEY = cfp.get("Contrast UI Details", "contrast.servicekey")
        except:
            print("Please specify contrast.servicekey in the script.properties")
        try:
            self.TEAMSERVER_URL = cfp.get("Contrast UI Details", "contrast.teamserver.url")
        except:
            print("Please specify contrast.teamserver.url in the script.properties")

        try:
            self.FOUND_DATE = cfp.get("Vulnerability Trend Configuration", "found.date")
        except:
            print("Please specify found.date in the script.properties")
        try:
            self.SORTING = cfp.get("Vulnerability Trend Configuration", "sorting.method")
        except:
            print("Please specify sorting.method in the script.properties")
        try:
            self.LICENSED_ONLY = cfp.get("Vulnerability Trend Configuration", "licensed.only") in ["True", 'true']
        except:
            print("Please specify licensed.only in the script.properties")

        try:
            self.startingMonth = int(cfp.get("Vulnerability Trend Duration", "starting.month"))
        except:
            print("Please specify starting.month in the script.properties")
        try:
            self.startingYear = int(cfp.get("Vulnerability Trend Duration", "starting.year"))
        except:
            print("Please specify starting.year in the script.properties")
        try:
            self.startingDay = int(cfp.get("Vulnerability Trend Duration", "starting.day"))
        except:
            print("Please specify starting.day in the script.properties")
        try:
            self.endingMonth = int(cfp.get("Vulnerability Trend Duration", "ending.month"))
        except:
            print("Please specify ending.month in the script.properties")
        try:
            self.endingDay = int(cfp.get("Vulnerability Trend Duration", "ending.day"))
        except:
            print("Please specify ending.day in the script.properties")
        try:
            self.endingYear = int(cfp.get("Vulnerability Trend Duration", "ending.year"))
        except:
            print("Please specify ending.year in the script.properties")

        self.AUTHORIZATION = base64.b64encode((self.USERNAME + ':' + self.SERVICE_KEY).encode('utf-8'))
        self.header = {
            "Authorization": self.AUTHORIZATION,
            "API-Key": self.API_KEY
        }
        self.getOrganizations()

    def getOrganizations(self):
        endpoint = '/organizations'
        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + endpoint

        response = requests.get(url=url, headers=self.header, stream=True)
        organization = json.loads(response.text)
        self.outputpath += '/' + organization['organization']['name']
        if not os.path.exists(self.outputpath):
            os.makedirs(self.outputpath)

    def getOfflineServers(self):
        """
        Output a list of all servers which are currently offline
        """

        endpoint = self.ORGANIZATION_ID + "/servers/filter?expand=applications,server_license," \
                                          "skip_links&includeArchived=false&offset=0&quickFilter=OFFLINE" \
                                          "&sort=-lastActivity"
        url = self.TEAMSERVER_URL + endpoint

        # Get response
        response = requests.get(url=url, headers=self.header, stream=True)
        jsonreader = json.loads(response.text)

        # Setup file to output results to
        filename = self.outputpath + "/OfflineServers.csv"
        filewriter = open(filename, 'w+')

        # Loop through each server and determine if it is offline
        if jsonreader["success"] is True:
            todaydate = datetime.today()
            servernum = 1
            print("The following servers are offline as of %s:" % todaydate)
            filewriter.write("The following servers are offline as of %s:\n" % todaydate)
            for server in jsonreader["servers"]:
                if server['status'] == "OFFLINE":  # If the status is offline, add it to our list
                    linetowrite = ("\t%d,%s" % (servernum, server['name']))
                    filewriter.write(linetowrite + "\n")
                    print(linetowrite)
                    servernum += 1
        filewriter.close()

    def getNeverLoggedInUsers(self):
        """
        Output a list of all users who have never logged into the Contrast UI
        """

        endpoint = self.ORGANIZATION_ID + "/users?expand=preferences,login,role," \
                                          "skip_links&offset=0&q=&quickFilter=ALL&sort=name"
        url = self.TEAMSERVER_URL + endpoint

        header = {
            "API-Key": self.API_KEY,
            "Authorization": self.AUTHORIZATION,
        }

        # Get response from the request
        response = requests.get(url=url, headers=header, stream=True)
        jsonreader = json.loads(response.text)

        # Setup file to output results to
        filename = self.outputpath + "/NeverLoggedInUsers.csv"
        filewriter = open(filename, 'w+')

        # Loop through each user to determine if they have ever logged in
        if jsonreader["success"] is True:
            user_num = 1
            print("\nThe following users have never logged into the teamserver")
            filewriter.write("The following users have never logged into the teamserver\n")
            for user in jsonreader["users"]:

                # Get their last login time. If there is last_login_time, they've never logged and we add it to the list
                try:
                    # noinspection PyUnusedLocal
                    lastlogintime = user['login']['last_login_time']
                except:
                    linetowrite = ("\t%d,%s" % (user_num, user['user_uid']))
                    filewriter.write(linetowrite + "\n")
                    print(linetowrite)
                    user_num += 1
        filewriter.close()

    def getUsersNotLoggedInDays(self, days):
        """
        Returns all users who have not logged in during the last "days" number of days
        Fail the command if the number of days is not positive

        :param days: number of days to check user login status for
        """

        if days < -1:
            print("\n** The number of days must be greater than 0")

        else:
            endpoint = self.ORGANIZATION_ID + "/users?expand=preferences,login,role," \
                                              "skip_links&offset=0&q=&quickFilter=ALL&sort=name"
            url = self.TEAMSERVER_URL + endpoint
            header = {
                "API-Key": self.API_KEY,
                "Authorization": self.AUTHORIZATION,
            }

            # Send the request and get its response
            response = requests.get(url=url, headers=header, stream=True)
            jsonreader = json.loads(response.text)

            # Setup file to output the text to
            filestring = "/UsersWhoHaventLoggedIn_%d_days.csv" % days
            filename = self.outputpath + filestring
            filewriter = open(filename, 'w+')

            # Loop through each user and get their last login time, compare it to the specified # of days
            if jsonreader["success"] is True:
                print("\nGetting users who have NOT logged into the teamserver in the past %d day(s)..." % days)
                filewriter.write(
                    "Username,Days Not Logged In\n")
                usercount = 1
                for user in jsonreader["users"]:
                    try:
                        # Get the user's last login time and convert it to a string
                        lastlogintime = datetime.fromtimestamp(
                            user['login']['last_login_time'] / 1000.0).strftime(
                            '%Y-%m-%d')

                        # Parse out month, day and year and find the difference between today and that date
                        year, month, day = lastlogintime.split("-")
                        # dt_lastlogintime = datetime(int(year), int(month), int(day))
                        date_diff = (datetime.today() - datetime(int(year), int(month), int(day))).days

                        # If difference between dates is greater than the one specified, add it to our output
                        if float(date_diff) > float(days):
                            linetowrite = (
                                "%s,%d" % (user['user_uid'], date_diff))
                            filewriter.write(linetowrite + "\n")
                            usercount += 1
                    except Exception:
                        linetowrite = (
                            "%s,%s" % (user['user_uid'], "Never_Logged_In"))
                        filewriter.write(linetowrite + "\n")
                        usercount += 1
                        continue
            filewriter.close()
            print("\t- Metrics written to %s" % filename)

    def getApplicationsWithNoGroup(self):
        """
        Output a list of all applications which do not have a group
        """

        endpoint = self.ORGANIZATION_ID + "/groups"
        url = self.TEAMSERVER_URL + endpoint

        header = {
            "API-Key": self.API_KEY,
            "Authorization": self.AUTHORIZATION,
        }

        # Get response from the request
        response = requests.get(url=url, headers=header, stream=True)
        jsonreader = json.loads(response.text)

        # Setup file to output results to
        filename = self.outputpath + "/ApplicationsWithNoGroup.csv"
        filewriter = open(filename, 'w+')

        if jsonreader["success"] is True:

            # Get all group ids from custom groups
            group_ids = []
            custom_groups = jsonreader['custom_groups']['groups']
            for custom_group in custom_groups:
                # print(custom_group['name'])
                group_ids.append(custom_group['group_id'])

            # Get all group ids from default groups
            predefined_groups = jsonreader['predefined_groups']['groups']
            for predefined_group in predefined_groups:
                group_ids.append(predefined_group['group_id'])

            # Query each group and build a list of all apps registered within a group
            registered_application_ids = []
            for group_id in group_ids:
                url = self.TEAMSERVER_URL + "/" + self.ORGANIZATION_ID + "/groups/" + str(group_id)
                response = requests.get(url=url, headers=header, stream=True)
                jsonreader = json.loads(response.text)

                if jsonreader["success"] is True:

                    # Make sure the group has an application. Loop through each group and get all app ids which are
                    # included in a group.
                    if jsonreader['group']['applications'] is not None:
                        applications_group = jsonreader['group']['applications']
                        for application_group in applications_group:
                            applications = application_group['applications']
                            for application in applications:
                                registered_application_ids.append(application['app_id'])
                else:
                    print("\nUnable to query individual group id")

            # Get a list of all applications. Need to exclude merged and archived applications.
            endpoint = self.ORGANIZATION_ID + "/applications?includeMerged=false&includeArchived=false"
            url = self.TEAMSERVER_URL + endpoint

            response = requests.get(url=url, headers=header, stream=True)
            jsonreader = json.loads(response.text)

            # Loop through each application and see if it's in our list of registered apps
            if jsonreader["success"] is True:
                applications = jsonreader['applications']
                missing_app_ids = {}

                # If the application is not in our list of registered apps, add it to our final list
                for application in applications:
                    if application['app_id'] not in registered_application_ids:
                        missing_app_ids[application['app_id']] = application['name']

                # Loop through all missing applications and output it to the text file
                print("\nThe following applications are not included in a group:")
                filewriter.write("The following applications are not included in a group:\n")
                application_num = 1
                for missing_app_id, missing_app_name in missing_app_ids.items():
                    linetowrite = ("\t%d. %s" % (application_num, missing_app_name))
                    filewriter.write(linetowrite + "\n")
                    print(linetowrite)
                    application_num += 1

            else:
                print("\nUnable to get the list of all applications")

        else:
            print("\nUnable to connect to the teamserver")
        filewriter.close()

    def getServersWithNoApplications(self):
        """
        Output a list of all servers which don't have any applications associated with it
        """

        # The endpoint automatically returns all servers with no applications
        endpoint = self.ORGANIZATION_ID + '/servers/filter?applicationsIds=None'
        url = self.TEAMSERVER_URL + endpoint
        header = {
            "API-Key": self.API_KEY,
            "Authorization": self.AUTHORIZATION,
        }

        response = requests.get(url=url, headers=header, stream=True)
        jsonreader = json.loads(response.text)

        filename = self.outputpath + "/ServersWithNoApplications.csv"
        filewriter = open(filename, 'w+')

        # Loop through all servers
        if jsonreader["success"] is True:
            servernumber = 1
            print("\nGetting servers which do not have an application")
            filewriter.write("Server Name,Server Hostname")

            # Add server to our list
            for server in jsonreader['servers']:
                linetowrite = (
                    "\t%s,%s" % (server['name'], server['hostname']))
                filewriter.write(linetowrite + "\n")
                servernumber += 1
            filewriter.close()
            print("\t- Metrics written to %s" % filename)

    def getUsersInGroups(self, return_object, applications=None):
        """
        Output a list of all users in each group
        :param return_object: format to return users and applications in. Choices are "string" or "list" with every user.
        :param applications: if specified, output a list of all users associated with each application.
        """
        endpoint = "/groups?expand=users,applications,skip_links&offset=0&q=&quickFilter=ALL&sort=name"
        url = self.TEAMSERVER_URL + "/" + self.ORGANIZATION_ID + "/" + endpoint

        header = {
            "API-Key": self.API_KEY,
            "Authorization": self.AUTHORIZATION,
        }

        # Send API request
        if not self.groups:
            response = requests.get(url=url, headers=header, stream=True)
            groups = json.loads(response.text)
            self.groups = groups

        # Specify output file based on whether a list of applications is provided
        if applications is None:
            filename = self.outputpath + "/UsersInGroups.csv"
            filewriter = open(filename, 'w+')
        else:
            filename = self.outputpath + "/UsersForApp.csv"
            filewriter = open(filename, 'w+')

        if self.groups['success'] is True:

            # Build up list of all group IDs
            group_ids = []
            for custom_group in self.groups['custom_groups']['groups']:
                group_ids.append(custom_group['group_id'])
            for predefined_group in self.groups['predefined_groups']['groups']:
                group_ids.append(predefined_group['group_id'])

            if applications is None:
                filewriter.write("Group Name,First Name,Last Name,Email Address\n")
                for group_id in group_ids:
                    url = self.TEAMSERVER_URL + "/" + self.ORGANIZATION_ID + "/groups/" + str(group_id)

                    # Send request
                    response = requests.get(url=url, headers=header, stream=True)
                    jsonreader = json.loads(response.text)
                    groups = jsonreader['group']

                    try:
                        for user in groups['users']:
                            filewriter.write(
                                groups['name'] + ',' + user['first_name'] + ',' + user['last_name'] + ',' + user[
                                    'uid'] + "\n")
                    except Exception as e:
                        print(e)
                        continue

            else:
                # Master dictionary of mappings between applications and users
                group_apps_users = {}

                for group_id in group_ids:

                    if group_id not in self.group_details.keys():
                        # Reach out to the group's endpoint and get a list of all applications and users
                        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/groups/" + str(group_id)

                        # Send request
                        response = requests.get(url=url, headers=header, stream=True)
                        groups = json.loads(response.text)['group']
                        self.group_details[group_id] = groups
                    else:
                        groups = self.group_details[group_id]

                    try:
                        users = []
                        applications_in_group = []
                        # Get all users in the group
                        for user in groups['users']:
                            users.append(user['uid'])

                        # Get all applications in the group. If None, the group has access to ALL applications
                        if groups['applications'] is None:
                            applications_in_group = None
                        else:
                            for application_in_group in groups['applications']:
                                for app in application_in_group['applications']:
                                    applications_in_group.append(app['name'])

                        # Add users and applications to the master mapping dictionary
                        group_apps_users[groups['name']] = {'users': users, 'apps': applications_in_group}
                    except Exception as e:
                        print(e)

                # Generate the spreadsheet of application <=> user mapping
                app_user_mappings_string = {}
                app_user_mappings_list = {}
                for search_app in applications:
                    app_user_mappings_list[search_app] = []
                    for groups, app_user in group_apps_users.items():
                        app_user_mapping = ""
                        apps = app_user['apps']
                        users = app_user['users']
                        users_list = []

                        # If ALL applications, add this application to every user
                        if apps is None:
                            for user in users:
                                app_user_mappings_list[search_app].append(user)
                                app_user_mapping += user + ','
                            if search_app in app_user_mappings_string.keys():
                                appname_without_comma = search_app.replace(",", "-")
                                app_user_mappings_string[appname_without_comma] += app_user_mapping
                            else:
                                appname_without_comma = search_app.replace(",", "-")
                                app_user_mappings_string[appname_without_comma] = app_user_mapping
                        else:
                            # Determine if the current user has access to one of the apps in the specified app list
                            if search_app in apps:
                                # app_user_mapping += search_app + ','
                                for user in users:
                                    app_user_mapping += user + ','
                                    app_user_mappings_list[search_app].append(user)
                                if search_app in app_user_mappings_string.keys():
                                    appname_without_comma = search_app.replace(",", "-")
                                    app_user_mappings_string[appname_without_comma] += app_user_mapping
                                else:
                                    appname_without_comma = search_app.replace(",", "-")
                                    app_user_mappings_string[appname_without_comma] = app_user_mapping

                # Generate csv of user and app mappings
                for app_name, users in app_user_mappings_string.items():
                    linetowrite = app_name + ',' + users + '\n'
                    filewriter.write(linetowrite)
                filewriter.close()
            # print("\t- Metrics output to %s" % filename)
            if return_object is "list":
                return app_user_mappings_list
            if return_object is "string":
                return app_user_mappings_string
            if return_object is "both":
                return app_user_mappings_string, app_user_mappings_list

    def getPercentUsersLoggedIn(self, days):
        """
        Determine percent of users who have not logged into the teamserver in the specified number of days
        :param days: number of days to check login status for
        """
        endpoint = self.ORGANIZATION_ID + "/users?expand=preferences,login,role," \
                                          "skip_links&offset=0&q=&quickFilter=ALL&sort=name"
        url = self.TEAMSERVER_URL + endpoint
        header = {
            "API-Key": self.API_KEY,
            "Authorization": self.AUTHORIZATION,
        }
        response = requests.get(url=url, headers=header, stream=True)
        jsonreader = json.loads(response.text)

        # Loop through each user and get their last login time, compare it to the specified # of days
        if jsonreader["success"] is True:
            total_users = jsonreader["users"].__len__()
            usercount = 0
            for user in jsonreader["users"]:
                try:
                    # Get the user's last login time and convert it to a string
                    lastlogintime = datetime.fromtimestamp(
                        user['login']['last_login_time'] / 1000.0).strftime(
                        '%Y-%m-%d')

                    # Parse out month, day and year and find the difference between today and that date
                    year, month, day = lastlogintime.split("-")
                    date_diff = (datetime.today() - datetime(int(year), int(month), int(day))).days

                    # If difference between dates is greater than the one specified, add it to our output
                    if float(date_diff) < float(days):
                        usercount += 1
                    else:
                        print(days)

                except Exception as e:
                    # print(e)
                    continue
            login_percentage = int(round(float(float(usercount) / float(total_users)) * 100.0))
            print(str(login_percentage) + "% of users have logged into teamserver the past " + str(days) + " days.\n")

    def parseAuditLog(self, days):
        """
        Parse the audit log and output a csv of some of the actions. The actions which will be logged are:
        - date
        - action: trace_status_changed, trace_deleted, agent_downloaded, license_applied
        - username
        - message
        :param days: Number of days to pull the audit log details for
        """

        # Setup file to output the text to
        filestring = "/usage_metrics.csv"
        filename = self.outputpath + filestring
        filewriter = open(filename, 'w+')

        # Generate epoch timestamp of earlier date
        earlier_date = (datetime.today() - timedelta(days)).strftime("%Y-%m-%d")
        endpoint = self.ORGANIZATION_ID + "/security/audit?expand=skip_links&limit=1000000&startDate=%s" % earlier_date

        url = self.TEAMSERVER_URL + endpoint
        header = {
            "API-Key": self.API_KEY,
            "Authorization": self.AUTHORIZATION,
        }

        # Send the request and get its response
        response = requests.get(url=url, headers=header, stream=True)
        jsonreader = json.loads(response.text)

        # Initialize keyword and counters for each action
        trace_status_change_counter = 0
        trace_num = 0

        trace_status_keyword = "trace_status_change"

        trace_deleted_counter = 0
        trace_deleted_keyword = "trace_deleted"

        report_usage_counter = 0
        report_usage_keyword = "report_generated"

        agent_downloaded_counter = 0
        agent_downloaded_keyword = "agent_download"

        license_applied_counter = 0
        license_applied_keyword = "license_applied"

        print("Parsing the audit log")
        print("\tTotal number of lines: ", jsonreader['logs'].__len__())

        # Setup headers for csv to output
        filewriter.write("date,action,username,message\n")
        for log in jsonreader['logs']:
            log_message = log['message']

            # Get epoch timestamp of audit message
            epoch_timestamp = log['date']
            timestamp = datetime.fromtimestamp(epoch_timestamp / 1000.0).strftime('%Y-%m-%d')
            username = ""
            message = ""
            keyword = ""
            line = ""

            # Check if audit message is for a trace status change, increment counter if so
            trace_status_matcher = re.compile(
                '([\w\.]+@[\w\.]+com)\smarked status to (\w+).+\s[for traces\(\)]+([\w\d\-\,]+)')
            trace_status_matches = trace_status_matcher.findall(log_message)
            if trace_status_matches:
                username = trace_status_matches[0][0]
                message = trace_status_matches[0][1]
                keyword = trace_status_keyword
                trace_status_change_counter += 1
                trace_num = trace_status_matches[0][2].count(",") + 1

            # Check if audit message is for a deleted trace, increment counter if so
            elif log_message.__contains__("deleted trace"):
                trace_deleted_matcher = re.compile("User\s([\.\@\w]+)\sdeleted\strace\s([A-Z0-9\-]+)")
                matches = trace_deleted_matcher.findall(log_message)
                try:
                    username = matches[0][0]
                    message = matches[0][1]
                except:
                    username = "na"
                    message = "na"
                keyword = trace_deleted_keyword
                trace_deleted_counter += 1

            # Check if audit message is for a report generation, increment counter if so
            elif log_message.__contains__("report"):
                report_downloaded_matcher = re.compile("User\s([\.\@\w]+)\screated a new report")
                matches = report_downloaded_matcher.findall(log_message)
                try:
                    username = matches[0]
                except:
                    username = ""
                message = "report generated"
                keyword = report_usage_keyword
                report_usage_counter += 1

            # Check if audit message is for an agent download, increment counter if so
            elif log_message.__contains__("downloaded") and not log_message.__contains__("agent_"):
                agent_download_matcher = re.compile('User\s([\w\@\.]+com)[downloaded\s]+([\w\s\_\.]+)[a|A]gent')
                keyword = agent_downloaded_keyword
                try:
                    matches = agent_download_matcher.findall(log_message)
                    username = matches[0][0]
                    if log_message.__contains__("JAVA_LAUNCHER"):
                        message = "JAVA_LAUNCHER"
                    else:
                        message = matches[0][1]
                except:
                    message = "na"
                    username = "na"
                agent_downloaded_counter += 1

            # Check if audit message is for a license applied, increment counter if so
            elif log_message.__contains__("License"):
                license_applied_matcher = re.compile(
                    "Enterprise License applied to application:[\s\\\']+(.+)\'[\\\'\s]+by\s(.+)")
                license_applied_matches = license_applied_matcher.findall(log_message)
                keyword = license_applied_keyword
                try:
                    username = license_applied_matches[0][1]
                    message = license_applied_matches[0][0]
                except:
                    username = "na"
                    message = "na"
                license_applied_counter += 1

            # Write line to csv
            try:
                if username is not "":
                    if keyword is "trace_status_change":
                        for trace in range(0, trace_num):
                            line = timestamp + ',' + keyword + ',' + username + ',' + message + ',\n'
                            filewriter.write(line)
                    else:
                        line = timestamp + ',' + keyword + ',' + username + ',' + message + ',\n'
                        filewriter.write(line)
            except:
                continue

        filewriter.close()

    def UsageMetrics(self, days):

        # Fail the command if the number of days is not positive
        """
        Get metrics related to usage of the Contrast UI.
        - Percent of users who have logged in within X days
        - Parse audit log to determine actions taken in the Contrast UI
        :param days:
        """
        if days < -1:
            print("\n* * The number of days must be greater than 0")

        else:
            self.getPercentUsersLoggedIn(days)
            self.parseAuditLog(days)

    def getDateRange(self):
        """
        Get date range from user if not specified in properties at the top of the script
        """
        if self.startingMonth is 0:
            print("Please specify the date range to pull vulnerability information.")
            self.startingMonth = int(input("\tStarting month number (1-12): "))
            while self.startingMonth < 1 or self.startingMonth > 12:
                print("The month should be between 1 and 12!")
                self.startingMonth = int(input("\tStarting month number (1-12): "))

        if self.startingDay is 0:
            self.startingDay = int(input("\tStarting day (1-31): "))
            while self.startingDay < 1 or self.startingDay > 31:
                print("The day should be between 1 and 31!")
                self.startingDay = int(input("\tStarting day (1-31): "))

        if self.endingMonth is 0:
            self.endingMonth = int(input("\tEnding month number (1-12): "))
            while self.endingMonth < 1 or self.endingMonth > 12:
                print("The month should be between 1 and 12!")
                self.endingMonth = int(input("\tEnding month number (1-12): "))

        if self.endingDay is 0:
            self.endingDay = int(input("\tEnding day (1-31): "))
            while self.endingDay < 1 or self.endingDay > 31:
                print("The day should be between 1 and 31!")
                self.endingDay = int(input("\tEnding day (1-31): "))

    def VulnerabilityTrendMetrics(self, application_metrics, printMetrics=True):
        # Get vulnerabilities per month for date range specified
        """
        Generate and output a csv of trending metrics for vulnerabilities in the Contrast UI.
        Date range is taken from properties defined at the top of the script.
        By default, it will query for ALL vulnerabilities (Open and Closed) and sort by the First Found date (as
            opposed to the Last Found date).
        Metrics are grouped by the month they were found in.

        - "Serious" vulnerabilities = vulnerabilities marked as Critical and High
        """

        print("\n\nGetting Vulnerability Metrics")
        yearlyMetrics = self.dateTrendManager_Organization(printMetrics=printMetrics)

        # Generate cumulative counts for all retrieved vulnerabilities
        cumulative_yearly_metrics, serious_categories = self.getCumulativeCounts(yearlyMetrics,
                                                                                 printMetrics=printMetrics)

        # Output the cumulative metrics to a file
        self.writeCumulativeMetrics(cumulative_yearly_metrics, serious_categories,
                                    application_metrics=application_metrics, printMetrics=printMetrics)

    def dateTrendManager_Organization(self, printMetrics=True):
        """
        Manages how vulnerabilities are associated with the time they were found
        :param printMetrics:
        :return:
        """

        # Get the months and days we'll be looking through
        self.getDateRange()

        yearlyMetrics = {}

        # Loop through all years
        year_index = self.startingYear

        # Check if metrics should be pulled for licensed applications only.
        # If true, generate a list of licensed application IDs
        if self.LICENSED_ONLY:
            applications = self.getApplications()
            self.application_licensed_status = self.getApplicationLicenseStatus(applications)

            # self.startTimeEpoch = int(
            #     datetime.datetime(year_index, month_index, 1, 0, 0, 0, 0).timestamp()) * 1000
            # self.endTimeEpoch = int(
            #     datetime.datetime(year_index, month_index, endingDay, 23, 59, 59, 99).timestamp()) * 1000

        # If time range spans multiple years, loop through all months in those years except the ending year
        if year_index < self.endingYear:
            # Loop through all months in  previous years
            starting_month_index = self.startingMonth
            while year_index < self.endingYear:
                monthlyMetrics = {}
                for month_index in range(starting_month_index, 13):
                    month = datetime(year_index, month_index, 1).strftime("%B")
                    endingDay = calendar.monthrange(year_index, month_index)[1]  # Get the number of days in the month
                    self.startTimeEpoch = int((
                                                  datetime(year_index, month_index, 1, 0, 0, 0, 0) - timedelta(
                                                      hours=3)).timestamp()) * 1000
                    self.endTimeEpoch = int((
                                                datetime(year_index, month_index, endingDay, 23, 59, 59,
                                                         99) - timedelta(hours=3)).timestamp()) * 1000
                    print("\n==========> Getting vulns in between %s %d, %d and %s %d, %d" % (
                        month, 1, year_index, month, endingDay, year_index))
                    monthlyMetrics[month] = self.getVulnsByDate()  # Get vulnerabilities for the month
                yearlyMetrics[year_index] = monthlyMetrics
                year_index += 1
                starting_month_index = 1

            # Loop through all months in the current year except the last month
            monthlyMetrics = {}
            starting_month_index = 1
            for month_index in range(starting_month_index, self.endingMonth):
                month = datetime(year_index, month_index, 1).strftime("%B")
                endingDay = calendar.monthrange(year_index, month_index)[1]  # Get the number of days in the month
                self.startTimeEpoch = int((
                                              datetime(year_index, month_index, 1, 0, 0, 0, 0) - timedelta(
                                                  hours=3)).timestamp()) * 1000
                self.endTimeEpoch = int((
                                            datetime(year_index, month_index, endingDay, 23, 59, 59, 99) - timedelta(
                                                hours=3)).timestamp()) * 1000
                print("\n==========> Getting vulns in between %s %d, %d and %s %d, %d" % (
                    month, 1, year_index, month, endingDay, year_index))
                monthlyMetrics[month] = self.getVulnsByDate()  # Get vulnerabilities for current month

            # Get vulns for the last month
            month_index = self.endingMonth
            month = datetime(year_index, month_index, 1).strftime("%B")
            self.startTimeEpoch = int((
                                          datetime(year_index, month_index, 1, 0, 0, 0, 0) - timedelta(
                                              hours=3)).timestamp()) * 1000
            self.endTimeEpoch = int((
                                        datetime(year_index, month_index, self.endingDay, 23, 59, 59, 99) - timedelta(
                                            hours=3)).timestamp()) * 1000
            print("\n==========> Getting vulns in between %s %d, %d and %s %d, %d" % (
                month, 1, year_index, month, self.endingDay, year_index))
            monthlyMetrics[month] = self.getVulnsByDate()  # Get vulnerabilities for current month
            yearlyMetrics[year_index] = monthlyMetrics

        else:

            # If the starting year and ending year are the same, loop through all months except the last month
            monthlyMetrics = {}
            for month_index in range(self.startingMonth, self.endingMonth):
                month = datetime(year_index, month_index, 1).strftime("%B")
                endingDay = calendar.monthrange(year_index, month_index)[1]  # Get the number of days in the month
                self.startTimeEpoch = int((
                                              datetime(year_index, month_index, 1, 0, 0, 0, 0) - timedelta(
                                                  hours=3)).timestamp()) * 1000
                self.endTimeEpoch = int((
                                            datetime(year_index, month_index, endingDay, 23, 59, 59, 99) - timedelta(
                                                hours=3)).timestamp()) * 1000
                print("\n==========> Getting vulns in between %s %d, %d and %s %d, %d" % (
                    month, 1, year_index, month, endingDay, year_index))
                monthlyMetrics[month] = self.getVulnsByDate()  # Get vulnerabilities for current month

            # Get vuln metrics for the last month
            month_index = self.endingMonth
            month = datetime(year_index, month_index, 1).strftime("%B")
            self.startTimeEpoch = int((
                                          datetime(year_index, month_index, 1, 0, 0, 0, 0) - timedelta(
                                              hours=3)).timestamp()) * 1000
            self.endTimeEpoch = int((
                                        datetime(year_index, month_index, self.endingDay, 23, 59, 59, 99) - timedelta(
                                            hours=3)).timestamp()) * 1000
            print("\n==========> Getting vulns in between %s %d, %d and %s %d, %d" % (
                month, 1, year_index, month, self.endingDay, year_index))
            monthlyMetrics[month] = self.getVulnsByDate()  # Get vulnerabilities for current month
            yearlyMetrics[year_index] = monthlyMetrics
        return yearlyMetrics

    def writeCumulativeMetrics(self, cumulativeMetrics, serious_categories, application_metrics, printMetrics):
        cumulative_metrics_filename = self.outputpath + "/CumulativeMetrics.csv"
        cumulative_filewriter = open(cumulative_metrics_filename, 'w+')

        serious_metrics = self.outputpath + "/SeriousMetrics.csv"
        serious_filewriter = open(serious_metrics, 'w+')

        if not application_metrics:
            metrics_linetowrite = [
                "Year,Month,Total Traces,Serious Traces,Cumulative Total Traces,Cumulative Serious Traces"]

            serious_category_header = "Year,Month"
            serious_category_total = 0
            serious_category_lines = []

            for category in serious_categories.keys():
                serious_category_header += ',' + category
            serious_category_header += ',Total'

            for year, monthlymetrics in cumulativeMetrics.items():
                for month, metrics in monthlymetrics.items():
                    monthly_serious_category_count = 0
                    serious_category_linetowrite = str(year) + ',' + str(month) + ','
                    mec_linetowrite = str(year) + ',' + str(month) + ',' + str(metrics['total_traces']) + ',' + str(
                        metrics[
                            'serious_traces']) + ',' + str(metrics['cumulative_total_traces']) + ',' + str(
                        metrics['cumulative_serious_traces'])
                    try:
                        # Loop through all serious categories found during date range
                        if serious_categories.__len__() > 0:
                            # Loop through all serious categories in the month
                            for category, total_category_count in serious_categories.items():
                                # try:
                                # If the category is the current serious metric
                                if category in metrics['serious_category_counts'].keys():
                                    count = metrics['serious_category_counts'][category]
                                    monthly_serious_category_count += count
                                    serious_category_linetowrite += str(count) + ','
                                else:
                                    serious_category_linetowrite += '0,'
                        serious_category_total += monthly_serious_category_count
                    except Exception as e:
                        print(e)

                    serious_category_linetowrite += str(monthly_serious_category_count)
                    metrics_linetowrite.append(mec_linetowrite)
                    serious_category_lines.append(serious_category_linetowrite)

            for line in metrics_linetowrite:
                cumulative_filewriter.write(line + '\n')

            serious_filewriter.write(serious_category_header + '\n')
            for line in serious_category_lines:
                serious_filewriter.write(line + '\n')
            linetowrite = "Total,Total"
            for count in serious_categories.values():
                linetowrite += ',' + str(count)
            linetowrite += ',' + str(serious_category_total)
            serious_filewriter.write(linetowrite)

        if application_metrics:
            metrics_linetowrite = [
                "Year,Month,Application Name,Environment,Critical,High,Medium,Low,Note,Total Traces,Serious Traces,"
                "Cumulative Total Traces,Cumulative Serious Traces"]

            serious_category_header = "Year,Month"
            serious_category_total = 0
            serious_category_lines = []
            # for category in serious_categories:
            #     serious_category_linetowrite += str(category) + ','

            for category in serious_categories.keys():
                serious_category_header += ',' + category
            serious_category_header += ',Total'
            for year, monthlymetrics in cumulativeMetrics.items():
                for month, metrics in monthlymetrics.items():
                    applications = metrics['applications']
                    for application, app_metrics in applications.items():
                        environments = app_metrics['environment']
                        for environment, environment_metrics in environments.items():
                            mec_linetowrite = str(year) + ',' + str(month) + ',' + application + ',' + environment + ',' \
                                              + str(environment_metrics['critical']) + ',' + str(
                                environment_metrics['high']) + ',' \
                                              + str(environment_metrics['medium']) + ',' + str(
                                environment_metrics['low']) + ',' \
                                              + str(environment_metrics['note']) + ',' + str(
                                metrics['total_traces']) + ',' + \
                                              str(metrics['serious_traces']) + ',' + str(
                                metrics['cumulative_total_traces']) + ',' \
                                              + str(metrics['cumulative_serious_traces'])
                            metrics_linetowrite.append(mec_linetowrite)

                    monthly_serious_category_count = 0
                    serious_category_linetowrite = str(year) + ',' + str(month) + ','

                    try:
                        # Loop through all serious categories found during date range
                        if serious_categories.__len__() > 0:
                            # Loop through all serious categories in the month
                            for category, total_category_count in serious_categories.items():
                                # try:
                                # If the category is the current serious metric
                                if category in metrics['serious_category_counts'].keys():
                                    count = metrics['serious_category_counts'][category]
                                    monthly_serious_category_count += count
                                    serious_category_linetowrite += str(count) + ','
                                else:
                                    serious_category_linetowrite += '0,'
                        serious_category_total += monthly_serious_category_count
                        # except:
                        #     serious_category_linetowrite += '0,'
                    except Exception as e:
                        print(e)

                    serious_category_linetowrite += str(monthly_serious_category_count)

                    serious_category_lines.append(serious_category_linetowrite)

            for line in metrics_linetowrite:
                cumulative_filewriter.write(line + '\n')

            serious_filewriter.write(serious_category_header + '\n')
            for line in serious_category_lines:
                serious_filewriter.write(line + '\n')
            linetowrite = "Total,Total"
            for count in serious_categories.values():
                linetowrite += ',' + str(count)
            linetowrite += ',' + str(serious_category_total)
            serious_filewriter.write(linetowrite)

        cumulative_filewriter.close()
        serious_filewriter.close()

    # Calculate cumulative vulnerability metrics
    def getCumulativeCounts(self, yearlymetrics, printMetrics=True):
        cumulative_total_counts = 0
        cumulative_serious_total_counts = 0
        serious_categories = {}
        for year, monthlymetrics in yearlymetrics.items():
            for month, metrics in monthlymetrics.items():
                try:
                    cumulative_total_counts += metrics['total_traces']
                    cumulative_serious_total_counts += metrics['serious_traces']

                    metrics['cumulative_total_traces'] = cumulative_total_counts
                    metrics['cumulative_serious_traces'] = cumulative_serious_total_counts

                    if metrics['serious_category_counts'].__len__() > 0:
                        monthly_serious_categories = metrics['serious_category_counts']
                        for category, count in monthly_serious_categories.items():
                            if category in serious_categories.keys():
                                serious_categories[category] += count
                            else:
                                serious_categories[category] = count
                    else:
                        pass
                    if printMetrics:
                        print(year, month, metrics['total_traces'], cumulative_total_counts, metrics['serious_traces'],
                              cumulative_serious_total_counts)
                except Exception as e:
                    print(e)
                    continue
        return yearlymetrics, serious_categories

    # Get vulnerabilities based on the specified date range
    def getVulnsByDate(self):
        print(
            "Fetch Vulnerabilities, sorting by " + self.SORTING + " and " + self.FOUND_DATE +
            " found in the date range...")

        try:
            # Get all vulns which need an issue opened for them
            limit = 200
            endpoint = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/orgtraces" \
                                                                    "/filter/?endDate=" + str(
                self.endTimeEpoch) + "&expand=application,servers,violations,bugtracker," \
                                     "skip_links&quickFilter=" + self.SORTING + \
                       ("&limit=%d&sort=-lastTimeSeen&startDate=" % limit) + str(self.startTimeEpoch) + \
                       "&timestampFilter=" + self.FOUND_DATE

            r = requests.get(url=endpoint, headers=self.header)
            vulns = json.loads(r.text)
            vuln_count = vulns['count']
            print("\t- Number of vulns:", vuln_count)

            for offset in range(limit, vuln_count, limit):

                endpoint = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/orgtraces/filter/?endDate=" + \
                           str(self.endTimeEpoch) + \
                           "&expand=application,servers,violations,bugtracker,skip_links&quickFilter=" + \
                           self.SORTING + ("&limit=%d&offset=%d" % (limit, offset)) + \
                           "&sort=-lastTimeSeen&startDate=" + str(self.startTimeEpoch) + \
                           "&timestampFilter=" + self.FOUND_DATE

                r = requests.get(url=endpoint, headers=self.header)
                next_vulns = json.loads(r.text)
                for vuln in next_vulns['traces']:
                    vulns['traces'].append(vuln)
                print("\t\tVulns picked up: ", offset)
            print("\t\tVulns picked up: ", vulns['traces'].__len__())
            return self.getVulnMetrics(vulns['traces'])

        except Exception as e:
            print("\n!! ERROR: Unable to connect to teamserver. Please check your authentication details.")
            print(e)
            return {}

    # Parse vulnerabilities and build metrics
    def getVulnMetrics(self, vulns):
        print("\nParsing Vulnerabilities...")

        traceNum = vulns.__len__()

        # total_traces = Total number of traces found in this time span
        # serious_traces = Total number of CRITICAL or HIGH vulns
        # serious_category_counts = Total number of serious categories and the number of vulns found in each one
        # changed_status = Number of vulns with a status other than Reported
        # remediated_status = Number of vulns with a status of Remediated, Not a Problem or Fixed
        metrics = {
            'serious_traces': 0,
            'changed_status': 0,
            'remediated_status': 0,
            'serious_category_counts': {},
            'applications': {},
            'total_traces': 0
        }

        application_traces = {}
        if traceNum > 0:
            print("\tTotal number of traces: %d" % traceNum)
            metrics['total_traces'] = traceNum
            seriousVulnCounter = 0
            seriousCategoryCounter = {}

            remediatedStatusCount = 0
            statusCount = 0

            unlicensedCounter = 0

            for vuln in vulns:
                if self.LICENSED_ONLY:
                    if vuln['application']['parent_app_id'] is not None:
                        vuln_app_id = vuln['application']['parent_app_id']
                    else:
                        vuln_app_id = vuln['application']['app_id']
                    if self.application_licensed_status[vuln_app_id] is False:
                        unlicensedCounter += 1
                    else:
                        # Check severity for serious vulns and increment counter for that severity
                        if vuln['default_severity'] in ("CRITICAL", "HIGH"):
                            seriousVulnCounter += 1
                            if vuln['rule_name'] in seriousCategoryCounter.keys():
                                seriousCategoryCounter[vuln['rule_name']] += 1
                            else:
                                seriousCategoryCounter[vuln['rule_name']] = 1

                        if vuln['status'] in ("Remediated", "Fixed", "Not a Problem", "Suspicious", "Confirmed"):
                            if vuln['status'] in ("Remediated", "Fixed", "Not a Problem"):
                                remediatedStatusCount += 1
                            statusCount += 1
                else:
                    if vuln['default_severity'] in ("CRITICAL", "HIGH"):
                        seriousVulnCounter += 1

                        if vuln['rule_name'] in seriousCategoryCounter.keys():
                            seriousCategoryCounter[vuln['rule_name']] += 1
                        else:
                            seriousCategoryCounter[vuln['rule_name']] = 1

                    if vuln['status'] in ("Remediated", "Fixed", "Not a Problem", "Suspicious", "Confirmed"):
                        if vuln['status'] in ("Remediated", "Fixed", "Not a Problem"):
                            remediatedStatusCount += 1
                        statusCount += 1

                    if vuln['application']['name'] in application_traces.keys():
                        application_traces[vuln['application']['name']]['total_traces'] += 1
                    else:
                        application_traces[vuln['application']['name']] = {'total_traces': 1}
                        application_traces[vuln['application']['name']].update(
                            {
                                'environment':
                                    {
                                        'QA': {
                                            'total_traces': 0,
                                            'critical': 0,
                                            'high': 0,
                                            'medium': 0,
                                            'low': 0,
                                            'note': 0
                                        },
                                        'DEVELOPMENT': {
                                            'total_traces': 0,
                                            'critical': 0,
                                            'high': 0,
                                            'medium': 0,
                                            'low': 0,
                                            'note': 0
                                        },
                                        'PRODUCTION': {
                                            'total_traces': 0,
                                            'critical': 0,
                                            'high': 0,
                                            'medium': 0,
                                            'low': 0,
                                            'note': 0
                                        }
                                    }
                            }
                        )
                    for server in vuln['servers']:
                        application_traces[vuln['application']['name']]['environment'][server['environment']][
                            'total_traces'] += 1
                        if vuln['severity'] == 'Critical':
                            application_traces[vuln['application']['name']]['environment'][server['environment']][
                                'critical'] += 1
                        if vuln['severity'] == 'High':
                            application_traces[vuln['application']['name']]['environment'][server['environment']][
                                'high'] += 1
                        if vuln['severity'] == 'Medium':
                            application_traces[vuln['application']['name']]['environment'][server['environment']][
                                'medium'] += 1
                        if vuln['severity'] == 'Low':
                            application_traces[vuln['application']['name']]['environment'][server['environment']][
                                'low'] += 1
                        if vuln['severity'] == 'Note':
                            application_traces[vuln['application']['name']]['environment'][server['environment']][
                                'note'] += 1

            print("\tTotal Unlicensed Vulnerabilities: ", unlicensedCounter)
            print("\n\tSerious Vulnerabilities Metrics")
            print("\t\t- Number of serious vulnerabilities: %d" % seriousVulnCounter)
            metrics['serious_traces'] = seriousVulnCounter

            category_counts = {}
            for category, count in seriousCategoryCounter.items():
                print("\t\t\t* %s: %d" % (category, count))
                category_counts[category] = count
            metrics['serious_category_counts'] = category_counts

            print("\n\tVulnerability Status")
            print("\t\t- Number of vulns in a status besides Reported:", statusCount)
            metrics['changed_status'] = statusCount

            print("\t\t- Number of vulnerabilities marked Remediated, Not a Problem or Fixed:",
                  remediatedStatusCount)
            metrics['remediated_status'] = remediatedStatusCount

            metrics['applications'] = application_traces

            return metrics

        else:
            print("\tThere were no traces found during the specified date range")
            metrics['serious_traces'] = 0
            metrics['changed_status'] = 0
            metrics['remediated_status'] = 0
            metrics['total_traces'] = 0
            metrics['serious_category_counts'] = {}
            return metrics

    def ApplicationMetrics(self):
        """
        Generate metrics for all applications in the teamserver.
        The following metrics will be logged in the csv.
        - Application name
        - # of Critical vulns
        - # of High Medium vulns
        - # of Medium vulns
        - # of Low vulns
        - # of Note vulns
        - # of vulns in a Reported status
        - # of vulns in a Suspicious status
        - # of vulns in a Confirmed status
        - # of vulns in a Remediated status
        - # of vulns in a Not a Problem status
        - # of vulns in a Fixed status
        - Application grade
        - First seen date for application
        - Last seen date for application
        - License level
        - Date license was applied
        """

        # Get all applications
        applications = self.getApplications()

        # Write metrics to specified log file
        self.writeApplicationMetrics(applications)

    def getUsersInTaggedApplications(self):
        untagged_applications = self.getApplicationsWithNoTag(search_tag_text="BU:")
        self.getUsersInGroups(applications=untagged_applications)

    # Get application metrics (count of applications)
    def getApplications(self):
        print("\nApplication Metrics")
        endpoint = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/applications/filter?expand=scores,license," \
                                                                "trace_breakdown,compliance_policy," \
                                                                "production_protected," \
                                                                "skip_links&filterText=&includeArchived=false" \
                                                                "&includeMerged=true&limit=100000&quickFilter" \
                                                                "=ALL&sort=appName"
        try:
            r = requests.get(url=endpoint, headers=self.header)
            if r.status_code == 504:
                limit = 50
                endpoint = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/applications/filter?expand=scores,license," \
                                                                        "trace_breakdown,compliance_policy," \
                                                                        "production_protected," \
                                                                        "skip_links&filterText=&includeArchived=false" \
                                                                        "&includeMerged=true&limit=50&quickFilter" \
                                                                        "=ALL&sort=appName"
                r = requests.get(url=endpoint, headers=self.header)
                applications = json.loads(r.text)
                app_count = applications['count']
                print("\t- Number of applications:", app_count)

                for offset in range(limit, app_count, limit):
                    endpoint = self.TEAMSERVER_URL + self.ORGANIZATION_ID + (
                        "/applications/filter?expand=scores,license," \
                        "trace_breakdown,modules,compliance_policy," \
                        "production_protected," \
                        "skip_links&filterText=&includeArchived=false" \
                        "&includeMerged=true&limit=%d&offset=%d"
                        "&quickFilter" \
                        "=ALL&sort=appName" % (limit, offset))
                    r = requests.get(url=endpoint, headers=self.header)
                    next_applications = json.loads(r.text)
                    for application in next_applications['applications']:
                        applications['applications'].append(application)
                    print("\t\tApplications picked up: ", offset)
                print("\t\tApplications picked up: ", applications['applications'].__len__())
                return applications
            else:
                applications = json.loads(r.text)
                print("\t- Number of applications:", applications['count'])
                return applications
        except:
            print("ERROR: Unable to retrieve applications")

    def writeApplicationMetrics(self, applications):

        filename = self.outputpath + '/ApplicationTraceBreakdown.csv'
        filewriter = open(filename, 'w+')

        app_severity_breakdown_header = "Application Name,Critical,High,Medium,Low,Note,Reported,Suspicious," \
                                        "Confirmed,Remediated,Not A Problem,Fixed,Grade,First Seen,Time Last Seen," \
                                        "License Level,Time Licensed\n "
        app_linetowrite = [app_severity_breakdown_header]

        for application in applications['applications']:
            app_severity_breakdown_linetowrite = ""
            appname = application['name'].replace(",", "-")
            app_linetowrite += appname + ','
            app_linetowrite += str(application['trace_breakdown']['criticals']) + ','
            app_linetowrite += str(application['trace_breakdown']['highs']) + ','
            app_linetowrite += str(application['trace_breakdown']['meds']) + ','
            app_linetowrite += str(application['trace_breakdown']['lows']) + ','
            app_linetowrite += str(application['trace_breakdown']['notes']) + ','
            app_linetowrite += str(application['trace_breakdown']['reported']) + ','
            app_linetowrite += str(application['trace_breakdown']['suspicious']) + ','
            app_linetowrite += str(application['trace_breakdown']['confirmed']) + ','
            app_linetowrite += str(application['trace_breakdown']['remediated']) + ','
            app_linetowrite += str(application['trace_breakdown']['notProblem']) + ','
            app_linetowrite += str(application['trace_breakdown']['fixed']) + ','
            app_linetowrite += str(application['scores']['letter_grade']) + ','
            app_linetowrite += datetime.fromtimestamp(application['created'] / 1000.0).strftime('%Y-%m-%d') + ','
            app_linetowrite += datetime.fromtimestamp(application['last_seen'] / 1000.0).strftime('%Y-%m-%d') + ','
            app_linetowrite += str(application['license']['level']) + ','
            if application['license']['start'] > 0:
                app_linetowrite += datetime.fromtimestamp(application['license']['start'] / 1000.0).strftime(
                    '%Y-%m-%d') + ','
            else:
                app_linetowrite += '0'
            app_linetowrite += '\n'
            app_linetowrite.append(app_severity_breakdown_linetowrite)

        for line in app_linetowrite:
            filewriter.write(line)

        filewriter.close()

    def getApplicationLicenseStatus(self, applications):
        app_license_status = {}
        if applications['count'] > 0:
            for application in applications['applications']:
                if application['license']['level'] == 'Licensed':
                    app_license_status[application['app_id']] = True
                else:
                    app_license_status[application['app_id']] = False
        return app_license_status

    def getApplicationLibraryMetricsByTag(self, search_tag_text):

        print("\nGetting vulnerable libraries for tagged applications\n\t- Searching for tag: '%s'" % search_tag_text)
        tags = self.getAllApplicationTags(filterText=search_tag_text)
        print("\t\t- Found %d tags which match filter text" % tags.__len__())

        taggedApplication_library_mappings = {}
        print("\t- Getting applications for each tag")
        for tag in tags:
            tagged_applications = self.getApplicationsWithTag(search_tag_text=tag)
            taggedApplication_library_mappings[tag] = tagged_applications

        print("\t- Getting vulnerable libraries for each application")
        for tag, apps in taggedApplication_library_mappings.items():
            for app in apps:
                app_id = app['app_id']
                vulnerable_libraries = self.getVulnerableLibraries_Application(app_id=app_id)
                app.update({"libraries": vulnerable_libraries})
                if vulnerable_libraries is not 0:
                    app_libraries = app['libraries']
                    for app_library in app_libraries:
                        library_vulns = self.getLibraryCVEs(library_hash=app_library['hash'],
                                                            library_language=app_library['app_language'])
                        app_library.update({'vulns': library_vulns})

        print(tagged_applications.__len__())

    def writeApplicationLibraryMetricsByTag(self, tagged_apps_vulnerable_libs):
        header = "Tag,Application Name,"

    def getAllApplicationTags(self, filterText=None):
        endpoint = "/applications/filters/tags/listing?expand=skip_links&filterText=&includeArchived=false&quickFilter=ALL"
        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + endpoint

        r = requests.get(url=url, headers=self.header)
        all_tags = json.loads(r.text)

        if filterText is not None:
            filter_tags = []
            for filter in all_tags['filters']:
                if filter['label'].find(filterText) is not -1:
                    filter_tags.append(filter['label'])
            return filter_tags
        else:
            return all_tags

    def getApplicationsWithNoTag(self, search_tag_text):
        # Text to search tags with
        endpoint = self.ORGANIZATION_ID + "/applications?includeMerged=false&includeArchived=false"
        url = self.TEAMSERVER_URL + endpoint

        response = requests.get(url=url, headers=self.header, stream=True)
        jsonreader = json.loads(response.text)

        # filename = '/appsWithNoBU.csv'
        # filewriter = open(filename, 'w+')

        applications = jsonreader['applications']
        app_count = 0
        untagged_applications = []
        for application in applications:
            tagged = False
            tags = application['tags']
            try:
                for tag in tags:
                    datatag = tag.find(search_tag_text)
                    if datatag is not -1:
                        tagged = True
            except:
                # print(application['name'])
                pass
            if not tagged:
                app_count += 1
                print(str(app_count) + ". " + application['name'])
                untagged_applications.append(application['name'])
        return untagged_applications

    def getApplicationsWithTag(self, search_tag_text):
        # Text to search tags with
        endpoint = self.ORGANIZATION_ID + "/applications?includeMerged=false&includeArchived=false"
        url = self.TEAMSERVER_URL + endpoint

        response = requests.get(url=url, headers=self.header, stream=True)
        jsonreader = json.loads(response.text)

        applications = jsonreader['applications']
        app_count = 0
        untagged_applications = []
        for application in applications:
            tagged = False
            tags = application['tags']
            try:
                for tag in tags:
                    datatag = tag.find(search_tag_text)
                    if datatag is not -1:
                        tagged = True
            except:
                # print(application['name'])
                pass
            if tagged:
                app_count += 1
                # print(str(app_count) + ". " + application['name'])
                untagged_applications.append({'name': application['name'], 'app_id': application['app_id']})
        return untagged_applications

    def getVulnerableLibraries_Application(self, app_id):
        endpoint = "/applications/%s/libraries/filter?expand=apps,quickFilters," \
                   "skip_links&q=&quickFilter=VULNERABLE&sort=score" % app_id
        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + endpoint

        r = requests.get(url=url, headers=self.header)
        app_vulnerable_libs = json.loads(r.text)
        if app_vulnerable_libs['count'] == 0:
            return 0
        else:
            return app_vulnerable_libs['libraries']

    def getLibraryCVEs(self, library_hash, library_language):
        print(library_hash)
        if library_language == "Java":
            url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/libraries/java/" + library_hash + \
                  "?expand=apps,vulns,skip_link "
        if library_language == ".NET":
            url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/libraries/dotnet/" + library_hash + \
                  "?expand=apps,vulns,skip_link "
        if library_language == "Node":
            url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/libraries/node/" + library_hash + \
                  "?expand=apps,vulns,skip_link "
        r = requests.get(url=url, headers=self.header)
        library_details = json.loads(r.text)
        return library_details['library']['vulns']

    def LibraryMetrics(self):
        print("\nGetting Library Metrics")
        # libraries = self.getLibraries()
        self.writeVulnerableLibraries(includeApplications=False)

    def getLibraries(self):
        endpoint = "/libraries"
        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + endpoint

        r = requests.get(url=url, headers=self.header, stream=True)
        libraries = json.loads(r.text)

        if libraries['success'] is True:
            print("\t- Found %d libraries" % libraries['libraries'].__len__())
            return libraries
        else:
            print("Unable to retrieve libraries. Please check the connection details.")
            return -1

    def writeVulnerableLibraries(self, includeApplications=True):
        endpoint = "/libraries/filter?expand=skip_links&q=&quickFilter=VULNERABLE&sort=score"
        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + endpoint

        r = requests.get(url=url, headers=self.header)
        vulnerable_libraries = json.loads(r.text)

        file_to_write = self.outputpath + "/VulnerableLibraries.csv"
        filewriter = open(file_to_write, 'w+')

        if includeApplications:
            file_header = "Language,Library Name,Months Outdated,Number of Known CVEs,Number of High Severity CVEs," \
                          "CVE IDs,Applications,Users\n"
        else:
            file_header = "Language,Library Name,Months Outdated,Number of Known CVEs,Number of High Severity CVEs," \
                          "CVE IDs\n"

        library_lines = [file_header]
        all_applications = {}

        if vulnerable_libraries['success'] is True:
            print("\t- Found %d vulnerable libraries! Parsing..." % vulnerable_libraries['count'])
            if vulnerable_libraries['count'] > 0:
                for count, library in enumerate(vulnerable_libraries['libraries']):
                    library_hash = library['hash']
                    if library['app_language'] == "Java":
                        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/libraries/java/" + library_hash + \
                              "?expand=apps,vulns,skip_link "
                    if library['app_language'] == ".NET":
                        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/libraries/dotnet/" + library_hash + \
                              "?expand=apps,vulns,skip_link "
                    if library['app_language'] == "Node":
                        url = self.TEAMSERVER_URL + self.ORGANIZATION_ID + "/libraries/node/" + library_hash + \
                              "?expand=apps,vulns,skip_link "
                    r = requests.get(url=url, headers=self.header)
                    library_details = json.loads(r.text)
                    if library_details['success'] is True:
                        line_to_write = library_details['library']['app_language'] + ',' + \
                                        library_details['library']['file_name'] + ',' + \
                                        str(library_details['library']['months_outdated']) + ',' + \
                                        str(library_details['library']['vulns'].__len__()) + ',' + \
                                        str(library_details['library']['high_vulnerabilities']) + ','
                        for vuln in library_details['library']['vulns']:
                            line_to_write += vuln['name'] + ' // '
                        line_to_write += ','

                        if includeApplications:
                            for app in library_details['library']['apps']:
                                line_to_write += app['name'] + " // "
                            line_to_write += ','

                            for app in library_details['library']['apps']:
                                if all_applications.__len__() > 0:
                                    if app['name'] not in all_applications.keys():
                                        users_list = self.getUsersInGroups(applications=[app['name']],
                                                                           return_object="list")
                                        all_applications[app['name']] = users_list[app['name']]
                                else:
                                    users_list = self.getUsersInGroups(applications=[app['name']], return_object="list")
                                    all_applications[app['name']] = users_list[app['name']]
                                users_in_app = all_applications[app['name']]
                                for user in users_in_app:
                                    line_to_write += user + ' // '
                            line_to_write += ','

                    line_to_write += "\n"
                    library_lines.append(line_to_write)

                    sys.stdout.write("\r\t\t%i libraries parsed" % (count + 1))
                    sys.stdout.flush()

        for line in library_lines:
            filewriter.write(line)
        filewriter.close()

    def test(self):
        a = datetime(2017, 12, 1, 0, 0, 0, 0)
        b = (a - timedelta(hours=3)).timestamp()
        print(a)
        print(b)

    def getTotalLicenses(self):

        endpoint = self.ORGANIZATION_ID + "/organizations/stats/licenses?expand=skip_links"

        url = self.TEAMSERVER_URL + endpoint
        try:
            response = requests.get(url=url, headers=self.header, stream=True)
            jsonreader = json.loads(response.text)
            # print (response.text)

            return_value = {}
            return_value['protect'] = jsonreader['total_protection']
            return_value['assess'] = jsonreader['total_assessment']

            # print (return_value)

            return return_value

        except Exception as e:
            print("ERROR: Unable to retrieve license info")
            print(e)
            print(response.text)

    def getLicenseHistory(self):

        endpoint = self.ORGANIZATION_ID + "/organizations/stats/licenses/history?expand=skip_links"

        url = self.TEAMSERVER_URL + endpoint
        try:
            response = requests.get(url=url, headers=self.header, stream=True)
            jsonreader = json.loads(response.text)
            # print (response.text)
            return jsonreader['license_history']
        except Exception as e:
            print("ERROR: Unable to retrieve license history")
            print(e)
            print(response.text)

    def writeLicenseHistory(self, license_history):

        licenses = self.getTotalLicenses()

        filename = self.outputpath + '/LicenseHistory.csv'
        filewriter = open(filename, 'w+')

        license_history_header = "Timestamp, Assess, Assess Available, Protect, Protect Available\n"
        license_history_lines = [license_history_header]

        for datapoint in license_history:

            if datapoint['assess'] is None and datapoint['protect'] is None:
                continue
            license_history_linetowrite = ""
            license_history_linetowrite += str(datetime.fromtimestamp(
                datapoint['timestamp'] / 1000.0).strftime(
                '%Y-%m-%d')) + ','
            # if licenses['assess'] != 0:
            license_history_linetowrite += str(datapoint['assess']) + ','
            license_history_linetowrite += str(licenses['assess']) + ','
            license_history_linetowrite += str(datapoint['protect']) + ','
            license_history_linetowrite += str(licenses['protect']) + '\n'
            license_history_lines.append(license_history_linetowrite)

        for line in license_history_lines:
            filewriter.write(line)

        filewriter.close()

    # Note this function requires: "pip3 install lxml" prior to generating graphs, also run
    # ApplicationMetrics() first
    def generatePPT(self):
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.chart import XL_LEGEND_POSITION
        from pptx.util import Inches
        prs = Presentation("CBRTemplate.pptx")

        print("Generating PPT")

        slide = prs.slides.add_slide(prs.slide_layouts[9])
        title_placeholder = slide.shapes.title
        title_placeholder.text = 'SERIOUS VULNERABILITIES BY APPLICATION'

        # define chart data ---------------------
        chart_data = ChartData()
        categories = []
        criticals = []
        highs = []
        with open(self.outputpath + '/ApplicationTraceBreakdown.csv') as csvDataFile:
            csvReader = csv.reader(csvDataFile)
            next(csvReader, None)
            for row in csvReader:
                # skip this app is there are no criticals or highs
                if row[1] != '0' or row[2] != '0':
                    categories.append(row[0])
                    criticals.append(row[1])
                    highs.append(row[2])

        chart_data.categories = categories
        chart_data.add_series('Critical Vulnerabilities', criticals)
        chart_data.add_series('High Vulnerabilities', highs)

        # add chart to slide --------------------
        x, y, cx, cy = Inches(2), Inches(1.05), Inches(10), Inches(6)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        self.writeLicenseHistory(self.getLicenseHistory())

        slide = prs.slides.add_slide(prs.slide_layouts[9])
        title_placeholder = slide.shapes.title
        title_placeholder.text = 'LICENSE ADOPTION'

        dates = []
        assess = []
        protect = []
        assess_available = []
        protect_available = []

        with open(self.outputpath + '/LicenseHistory.csv') as csvDataFile:
            csvReader = csv.reader(csvDataFile)
            next(csvReader, None)
            for row in csvReader:
                dates.append(row[0])
                assess.append(row[1])
                assess_available.append(row[2])
                protect.append(row[3])
                protect_available.append(row[4])

                # define chart data ---------------------
        chart_data = ChartData()
        chart_data.categories = dates
        chart_data.add_series('Assess Available', assess_available)
        chart_data.add_series('Protect Available', protect_available)
        chart_data.add_series('Assess', assess)
        chart_data.add_series('Protect', protect)
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(11.5), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.include_in_layout = False
        # chart.series[0].smooth = True

        prs.save(self.outputpath + '/CBRTemplate.pptx')


controller = controller()

##### General metrics
# controller.getServersWithNoApplications()
# controller.getUsersNotLoggedInDays(days=90)
# controller.getApplicationsWithNoGroup()
# controller.getNeverLoggedInUsers()
# controller.getOfflineServers()
# controller.getUsersInGroups()
# controller.UsageMetrics(days=365)
# controller.getPercentUsersLoggedIn(365)
# controller.getUsersInTaggedApplications()

##### Application metrics
# controller.ApplicationMetrics()

##### Vulnerability metrics
# controller.VulnerabilityTrendMetrics(application_metrics=True)

##### Library metrics
# controller.LibraryMetrics()
# controller.getApplicationLibraryMetricsByTag(search_tag_text="demo")

# Note this function requires: "pip3 install lxml" prior to generating graphs
# controller.generatePPT()

# controller.test()
